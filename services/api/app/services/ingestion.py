"""Document ingestion: parse PDF, PPTX, DOCX into structured text chunks.

Validates file size, normalizes unicode, cleans text, and logs all operations.
"""

import io
import logging
import re
import unicodedata

logger = logging.getLogger(__name__)

_MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB


def _clean_text(text: str) -> str:
    """Remove NULL bytes, normalize unicode, normalize whitespace, and strip."""
    # Remove NULL bytes
    text = text.replace("\x00", "").replace("\u0000", "")
    # Unicode NFC normalization
    text = unicodedata.normalize("NFC", text)
    # Normalize whitespace (collapse multiple spaces/tabs, keep newlines)
    text = re.sub(r"[^\S\n]+", " ", text)
    # Collapse multiple blank lines into at most two newlines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def parse_document(
    file_bytes: bytes, file_type: str, file_name: str
) -> list[dict]:
    """Parse a document and return a list of text chunks.

    Args:
        file_bytes: Raw bytes of the uploaded file.
        file_type: One of "pdf", "pptx", or "docx".
        file_name: Original file name (for logging).

    Returns:
        List of chunk dicts with keys:
            chunk_index, content, heading, chunk_type

    Raises:
        ValueError: If file_type is unsupported or file exceeds size limit.
    """
    if not file_bytes:
        raise ValueError("file_bytes cannot be empty")
    if not file_type:
        raise ValueError("file_type is required")
    if not file_name:
        raise ValueError("file_name is required")

    # File size check
    file_size = len(file_bytes)
    if file_size > _MAX_FILE_SIZE:
        raise ValueError(
            f"File '{file_name}' is {file_size / (1024 * 1024):.1f}MB, "
            f"exceeds maximum of 50MB"
        )

    parsers = {
        "pdf": _parse_pdf,
        "pptx": _parse_pptx,
        "docx": _parse_docx,
    }
    parser = parsers.get(file_type)
    if parser is None:
        raise ValueError(f"Unsupported file type: {file_type}")

    logger.info(
        "Parsing document: file_name=%s, file_type=%s, file_size=%.1fKB",
        file_name,
        file_type,
        file_size / 1024,
    )

    chunks = parser(file_bytes, file_name)

    logger.info(
        "Parsing complete: file_name=%s, chunks_generated=%d",
        file_name,
        len(chunks),
    )
    return chunks


def _parse_pdf(file_bytes: bytes, file_name: str) -> list[dict]:
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    total_pages = len(reader.pages)
    chunks: list[dict] = []
    skipped_count = 0

    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        text = _clean_text(text)
        if text:
            chunks.append(
                {
                    "chunk_index": len(chunks),
                    "content": text,
                    "heading": f"Page {i + 1}",
                    "chunk_type": "page",
                }
            )
        else:
            skipped_count += 1

    if skipped_count > 0:
        logger.warning(
            "PDF '%s': skipped %d empty pages out of %d total",
            file_name,
            skipped_count,
            total_pages,
        )

    logger.info(
        "PDF '%s': total_pages=%d, chunks_generated=%d, skipped=%d",
        file_name,
        total_pages,
        len(chunks),
        skipped_count,
    )
    return chunks


def _parse_pptx(file_bytes: bytes, file_name: str) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    total_slides = len(prs.slides)
    chunks: list[dict] = []
    skipped_count = 0

    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        heading: str | None = None

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = _clean_text(paragraph.text)
                    if para_text:
                        if heading is None:
                            heading = para_text
                        texts.append(para_text)

        content = "\n".join(texts)
        if content.strip():
            chunks.append(
                {
                    "chunk_index": len(chunks),
                    "content": content,
                    "heading": heading,
                    "chunk_type": "slide",
                }
            )
        else:
            skipped_count += 1

    if skipped_count > 0:
        logger.warning(
            "PPTX '%s': skipped %d empty slides out of %d total",
            file_name,
            skipped_count,
            total_slides,
        )

    logger.info(
        "PPTX '%s': total_slides=%d, chunks_generated=%d, skipped=%d",
        file_name,
        total_slides,
        len(chunks),
        skipped_count,
    )
    return chunks


def _parse_docx(file_bytes: bytes, file_name: str) -> list[dict]:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    chunks: list[dict] = []
    current_heading: str | None = None
    current_texts: list[str] = []

    def _flush() -> None:
        nonlocal current_heading, current_texts
        content = "\n".join(current_texts).strip()
        if content:
            chunks.append(
                {
                    "chunk_index": len(chunks),
                    "content": content,
                    "heading": current_heading,
                    "chunk_type": "section",
                }
            )
        current_heading = None
        current_texts = []

    for paragraph in doc.paragraphs:
        # Handle paragraph.style being None
        style_name = ""
        if paragraph.style and paragraph.style.name:
            style_name = paragraph.style.name

        if style_name.startswith("Heading"):
            _flush()
            heading_text = _clean_text(paragraph.text)
            current_heading = heading_text or None
        else:
            text = _clean_text(paragraph.text)
            if text:
                current_texts.append(text)

    _flush()

    logger.info(
        "DOCX '%s': total_paragraphs=%d, chunks_generated=%d",
        file_name,
        len(doc.paragraphs),
        len(chunks),
    )
    return chunks
